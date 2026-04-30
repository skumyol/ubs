#!/usr/bin/env python3
"""Directly patch the existing UBS pitch deck PPTX with current values."""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation

from src.analysis import narrative_shift_analysis


ROOT = Path("/Users/skumyol/Documents/GitHub/ubs")
DECK_PATH = ROOT / "deck" / "UBS_Pitch_Deck_AUTO.pptx"


def set_text(shape, text: str) -> None:
    if hasattr(shape, "text_frame"):
        shape.text_frame.clear()
        shape.text_frame.text = text
    else:
        shape.text = text


def update_table_from_df(table, df: pd.DataFrame) -> None:
    headers = list(df.columns)
    for c, header in enumerate(headers):
        table.cell(0, c).text = str(header)
    for r in range(1, min(len(table.rows), len(df) + 1)):
        row = df.iloc[r - 1]
        for c, col in enumerate(headers):
            table.cell(r, c).text = str(row[col])


def sanitize_peer_comps(df: pd.DataFrame) -> pd.DataFrame:
    cleaned = df.copy()
    cleaned = cleaned.replace({"nan": "N/A"})
    cleaned = cleaned.fillna("N/A")
    cleaned["sector"] = cleaned["sector"].replace(
        {"Power Equipment / Industrial Conglomerate": "Inverter / Storage Equipment"}
    )
    return cleaned


def main() -> None:
    prs = Presentation(str(DECK_PATH))

    pair = pd.read_csv(ROOT / "data/processed/valuation/pair_trade_summary.csv").iloc[0]
    trader = pd.read_csv(ROOT / "data/processed/valuation/trader_analysis.csv").iloc[0]
    backtest = pd.read_csv(ROOT / "data/processed/valuation/oil_sungrow_backtest.csv").iloc[0]
    signal_tracker = pd.read_csv(ROOT / "outputs/tables/signal_tracker_plan_format.csv")
    peer_comps = pd.read_csv(ROOT / "data/processed/valuation/peer_comps.csv")
    long_scenarios = pd.read_csv(ROOT / "data/processed/valuation/long_scenarios.csv")
    short_scenarios = pd.read_csv(ROOT / "data/processed/valuation/short_scenarios.csv")
    classified = pd.read_csv(ROOT / "data/processed/classified_paragraphs.csv")
    narrative = narrative_shift_analysis(classified)

    long_ret = float(pair["long_expected_return_pct"])
    short_ret = float(pair["short_expected_move_pct"])
    spread = float(pair["pair_spread_return_pct"])

    # Slide 2
    set_text(
        prs.slides[1].shapes[1],
        "\n".join(
            [
                "• Thesis: Own grid integration leader capturing State Grid RMB 4T capex + synchronous condenser breakthrough",
                f"• Long Dongfang Electric @ target upside: {long_ret:.1f}%",
                f"• Short Sungrow @ target downside: {short_ret:.1f}%",
                f"• Pair spread expected return: {spread:.1f}%",
                "• 3 pillars: (1) Structural electricity demand (2) Grid capex visibility (3) Earnings-quality divergence",
                "• Key catalysts: State Grid 4T RMB capex, Dongfang synchronous condenser orders, Sungrow margin scrutiny",
            ]
        ),
    )

    # Slide 4
    set_text(
        prs.slides[3].shapes[1],
        "\n".join(
            [
                "• China's 15th FYP shifts the bottleneck from generation buildout to system reliability and grid integration",
                "• State Grid plans RMB 4 trillion of fixed-asset investment across 2026-2030, a step-up versus the prior cycle",
                "• New-type power system priorities include source-grid-load-storage coordination and grid flexibility",
                "• Dongfang is directly exposed to synchronous condensers, transmission equipment, and grid-supporting power equipment",
                "• This makes China grid capex, not generic global clean-tech beta, the key macro driver for the pair",
            ]
        ),
    )

    # Slide 9
    slide9 = prs.slides[8]
    table9 = slide9.shapes[1].table
    scorecard_rows = pd.DataFrame(
        [
            ["Historical spread test", f"2Y simulated pair P&L was {backtest['pair_final_pnl_pct']:.1f}%; Dongfang Electric 2Y return {backtest['long_2y_return']:.1f}% vs Sungrow {backtest['short_2y_return']:.1f}%.", "Risk / contradiction"],
            ["Earnings durability", "Dongfang Electric delivered 2025 net profit growth of 31.1% on 12.8% revenue growth, backed by order visibility, while Sungrow moved from +22.0% full-year profit growth to -40.1% YoY in Q1 2026.", "Supports long Dongfang Electric / short Sungrow"],
            ["Valuation stretch", "Sungrow trades at a 8.9x P/E premium and 5.4x P/B gap to Dongfang Electric (45.0x/8.5x vs 36.1x/3.1x).", "Supports short Sungrow if margin recovery disappoints"],
            ["Technical setup", f"Sungrow is at {trader['technicals_sungrow_pct_52w_high']:.1f}% of its 52-week high with RSI {trader['technicals_sungrow_rsi']:.1f}, while Dongfang Electric is at {trader['technicals_dongfang_pct_52w_high']:.1f}% with RSI {trader['technicals_dongfang_rsi']:.1f}.", "Supports timing discipline"],
            ["Scenario valuation", f"Probability-weighted modeled pair spread is {spread:.1f}%.", "Forward prediction"],
        ],
        columns=["Pillar", "Evidence", "Signal"],
    )
    update_table_from_df(table9, scorecard_rows)

    # Slide 11
    slide11 = prs.slides[10]
    table11 = slide11.shapes[1].table
    tracker_df = signal_tracker[["signal_cluster", "Grid Equipment", "Inverter & Storage Equipment"]].copy()
    tracker_df.columns = ["Signal Cluster", "Grid Equipment", "Inverter & Storage Equipment"]
    update_table_from_df(table11, tracker_df)
    set_text(
        slide11.shapes[2],
        f"Thesis support score: {float(narrative['thesis_support_score']):.3f} | {narrative['interpretation']} | Grid signal share: {float(narrative['grid_signal_share'])*100:.0f}%",
    )

    # Slide 12
    slide12 = prs.slides[11]
    table12 = slide12.shapes[1].table
    peer_df = peer_comps[["company", "ticker", "sector", "p_e", "ev_ebitda", "roic", "revenue_growth", "profit_margin"]].copy()
    peer_df = sanitize_peer_comps(peer_df)
    peer_df.columns = ["Company", "Ticker", "Sector", "P/E", "EV/EBITDA", "ROIC", "Revenue Growth", "Profit Margin"]
    update_table_from_df(table12, peer_df)
    set_text(
        slide12.shapes[2],
        "Comps compare Dongfang's grid-backbone exposure against premium downstream inverter/storage valuation risk.",
    )

    # Slide 14 / 15
    update_table_from_df(prs.slides[13].shapes[1].table, long_scenarios)
    update_table_from_df(prs.slides[14].shapes[1].table, short_scenarios)

    # Slide 18
    set_text(
        prs.slides[17].shapes[1],
        "\n".join(
            [
                "• POSITION SIZING (Risk-Based):",
                "•   • Portfolio: $100M example | Max risk per trade: 2% ($2M)",
                f"•   • Pair volatility: {float(trader['volatilities_pair_vol']):.1f}% annual | Position size: {float(trader['position_sizing_recommended_position_pct']):.1f}% of portfolio (${float(trader['position_sizing_recommended_notional_mm']):.1f}M notional)",
                f"•   • Allocation: ${float(trader['position_sizing_long_dongfang_notional_mm']):.1f}M long Dongfang / ${float(trader['position_sizing_short_sungrow_notional_mm']):.1f}M short Sungrow (dollar-neutral)",
                "• ",
                "• CARRY COST (6-month hold):",
                "•   • Sungrow borrow cost: 2.5-7% annually",
                "•   • Sungrow dividend/carry checked before execution",
                "•   • Dongfang dividend (long receives): 1.5% → ~$9K income",
                "•   • Net carry: ~$13K-27K based on current trader analysis",
                "• ",
                "• LIQUIDITY & EXECUTION:",
                f"•   • Dongfang (HK): ~${float(trader['liquidity_dongfang_daily_dollar_volume_mm']):.0f}M daily volume | Execute in <1 day | HKEX access",
                "•   • Sungrow: highly liquid A-share leg; execute in slices after borrow is confirmed",
                "•   • Dongfang trades on HKEX (1072.HK) - no Stock Connect needed",
                "• ",
                "• TECHNICAL TIMING:",
                f"•   • Sungrow: RSI {float(trader['technicals_sungrow_rsi']):.1f} | {float(trader['technicals_sungrow_pct_52w_high']):.1f}% of 52-week high",
                f"•   • Dongfang: RSI {float(trader['technicals_dongfang_rsi']):.1f} | {float(trader['technicals_dongfang_pct_52w_high']):.1f}% of 52-week high",
            ]
        ),
    )

    # Slide 20
    set_text(
        prs.slides[19].shapes[2],
        f"Long Dongfang (+{long_ret:.0f}%) / Short Sungrow ({short_ret:.0f}%) = Pair spread +{spread:.0f}%",
    )

    prs.save(str(DECK_PATH))
    print(f'[SAVED] {DECK_PATH}')


if __name__ == "__main__":
    main()
