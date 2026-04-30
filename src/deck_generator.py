"""Auto-generate the UBS pitch deck as a PPTX file.

Assembles slides from analysis outputs, evidence pack, charts, and valuation.
Follows the 13-slide structure from base_plan.md.
"""

import json
import pandas as pd
from pathlib import Path
from typing import Dict, List, Optional

from src.pair_config import (
    DECK_SUBTITLE,
    DECK_TITLE,
    LONG_LEG,
    SHORT_LEG,
    SLIDE_TITLES,
    TEAM_LINE,
)

# Brand colors (UBS-ish professional palette)
COLOR_PRIMARY = "#E60028"        # UBS red
COLOR_SECONDARY = "#002F5D"      # UBS dark blue
COLOR_GRID = "#005F9E"           # Long leg color
COLOR_OIL = "#C97B00"            # Short leg color
COLOR_TEXT = "#1A1A1A"
COLOR_MUTED = "#666666"


def _install_pptx():
    """Ensure python-pptx is installed."""
    try:
        import pptx  # noqa
    except ImportError:
        import subprocess
        import sys
        subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])


def _hex_to_rgb(hex_color: str):
    """Convert hex string to RGBColor."""
    from pptx.dml.color import RGBColor
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_title_slide(prs, title: str, subtitle: str):
    """Add a title slide with UBS styling."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Color band
    from pptx.shapes.autoshape import Shape
    from pptx.enum.shapes import MSO_SHAPE
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.3)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = _hex_to_rgb(COLOR_PRIMARY)
    band.line.fill.background()

    # Title
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb(COLOR_SECONDARY)

    # Subtitle
    sb = slide.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(12), Inches(1.0))
    sf = sb.text_frame
    sf.word_wrap = True
    sp = sf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(22)
    sp.font.color.rgb = _hex_to_rgb(COLOR_MUTED)

    if title == DECK_TITLE:
        mb = slide.shapes.add_textbox(Inches(0.7), Inches(4.4), Inches(12), Inches(0.8))
        mf = mb.text_frame
        mp = mf.paragraphs[0]
        mp.text = TEAM_LINE
        mp.font.size = Pt(16)
        mp.font.color.rgb = _hex_to_rgb(COLOR_MUTED)

    return slide


def add_content_slide(prs, title: str, body_items: List[str]):
    """Add a standard content slide with title and bullet points."""
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb(COLOR_SECONDARY)

    # Body
    bb = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.5))
    bf = bb.text_frame
    bf.word_wrap = True

    for i, item in enumerate(body_items):
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = _hex_to_rgb(COLOR_TEXT)
        p.space_after = Pt(8)

    return slide


def add_chart_slide(prs, title: str, image_path: Path, caption: str = ""):
    """Add a slide with a single chart image."""
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb(COLOR_SECONDARY)

    # Image
    if image_path and image_path.exists():
        slide.shapes.add_picture(
            str(image_path), Inches(1.5), Inches(1.3),
            width=Inches(10), height=Inches(5.5)
        )

    # Caption
    if caption:
        cb = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.5), Inches(0.5))
        cp = cb.text_frame.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(12)
        cp.font.italic = True
        cp.font.color.rgb = _hex_to_rgb(COLOR_MUTED)

    return slide


def add_table_slide(prs, title: str, df: pd.DataFrame, caption: str = ""):
    """Add a slide with a DataFrame rendered as a table."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb(COLOR_SECONDARY)

    # Drop internal columns (prefix _)
    display_df = df.loc[:, [c for c in df.columns if not str(c).startswith("_")]].copy()

    n_rows, n_cols = len(display_df) + 1, len(display_df.columns)
    if n_rows == 1 or n_cols == 0:
        return slide

    # Size the table
    tbl_width = Inches(min(12, 1.8 * n_cols))
    tbl_height = Inches(min(5.5, 0.4 * n_rows + 0.4))
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.7), Inches(1.3),
        tbl_width, tbl_height,
    )
    tbl = tbl_shape.table

    # Header
    for j, col in enumerate(display_df.columns):
        cell = tbl.cell(0, j)
        cell.text = str(col).replace("_", " ").title()
        cell.fill.solid()
        cell.fill.fore_color.rgb = _hex_to_rgb(COLOR_SECONDARY)
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = _hex_to_rgb("#FFFFFF")

    # Rows
    for i, (_, row) in enumerate(display_df.iterrows(), start=1):
        for j, col in enumerate(display_df.columns):
            cell = tbl.cell(i, j)
            val = row[col]
            cell.text = "" if pd.isna(val) else str(val)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11)
                    run.font.color.rgb = _hex_to_rgb(COLOR_TEXT)

    # Caption
    if caption:
        cb = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.5), Inches(0.5))
        cp = cb.text_frame.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(12)
        cp.font.italic = True
        cp.font.color.rgb = _hex_to_rgb(COLOR_MUTED)

    return slide


def add_quote_slide(prs, title: str, quotes: List[Dict]):
    """Add a slide with supporting quotes."""
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = _hex_to_rgb(COLOR_SECONDARY)

    # Quotes
    bb = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.5))
    bf = bb.text_frame
    bf.word_wrap = True

    if not quotes:
        p = bf.paragraphs[0]
        p.text = "No supporting quotes available. Expand data collection."
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = _hex_to_rgb(COLOR_MUTED)
        return slide

    for i, q in enumerate(quotes[:3]):
        quote_text = q["quote"][:280]
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.text = f'"{quote_text}"'
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = _hex_to_rgb(COLOR_TEXT)
        p.space_after = Pt(4)

        p2 = bf.add_paragraph()
        p2.text = f"  — {q['source']}  |  {q['category']}  |  conf {q['confidence']}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = _hex_to_rgb(COLOR_MUTED)
        p2.space_after = Pt(12)

    return slide


def build_deck(
    classified_df: Optional[pd.DataFrame],
    signal_tracker_df: Optional[pd.DataFrame],
    evidence_pack: Optional[Dict],
    valuation_summary: Optional[Dict],
    peer_comps_df: Optional[pd.DataFrame],
    dcf_df: Optional[pd.DataFrame],
    long_scenarios_df: Optional[pd.DataFrame],
    short_scenarios_df: Optional[pd.DataFrame],
    charts_dir: Path,
    output_path: Path,
    narrative_shift: Optional[Dict] = None,
    include_per_slide_evidence: bool = False,
    consolidate_ai_slides: bool = True,
) -> Path:
    """Assemble the full UBS pitch deck.

    Saves to output_path and returns it.
    """
    _install_pptx()
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # -------- Slide 1: Executive Summary --------
    add_title_slide(
        prs,
        "Long the Grid, Short the Bottleneck",
        f"Hong Kong Track | Energy Transition | {DECK_SUBTITLE}",
    )

    # -------- Slide 2: One-Pager --------
    one_pager_items = [
        "Thesis: Own grid integration leader capturing State Grid RMB 4T capex + synchronous condenser breakthrough",
        f"Long Dongfang Electric @ target upside: "
        f"{valuation_summary['long_expected_return']}%" if valuation_summary else
        "Long Dongfang Electric - grid capex + 140B RMB backlog + tech breakthrough",
        f"Short {SHORT_LEG.name} @ target downside: "
        f"{valuation_summary['short_expected_return']}%" if valuation_summary else
        f"Short {SHORT_LEG.name} - premium clean-tech multiple + demand normalization risk",
        f"Pair spread expected return: "
        f"{valuation_summary['pair_spread_return']}%" if valuation_summary else
        "Pair spread captures earnings-quality differential",
        "3 pillars: (1) Structural electricity demand (2) Grid capex visibility (3) Earnings-quality divergence",
        f"Key catalysts: State Grid 4T RMB capex, Dongfang synchronous condenser orders, {SHORT_LEG.name} margin scrutiny",
    ]
    add_content_slide(prs, "Executive Summary", one_pager_items)

    # -------- Slide 3: Variant View --------
    variant_items = [
        "Consensus treats energy-transition equipment winners as interchangeable",
        "Variant view: energy security is no longer just securing barrels — it is reliable electricity",
        "Data centers, EVs, industrial electrification, and cooling all stress the power system",
        "Governments treat grid hardening as national security infrastructure",
        "→ Oil disruption is the symptom. Grid investment is the solution.",
    ]
    add_content_slide(prs, "Variant View: The Real Energy-Security Trade Is Grid Resilience", variant_items)

    # Slide 3/4 evidence (optional)
    if include_per_slide_evidence and evidence_pack and "slide_3_variant_view" in evidence_pack:
        add_quote_slide(
            prs,
            "Evidence: Grid Resilience Signals in Primary Sources",
            evidence_pack["slide_3_variant_view"]["quotes"],
        )

    # -------- Slide 4: Industry Outlook --------
    industry_items = [
        "Global electricity demand expected to grow ~4% annually through 2030 (IEA)",
        "US data center power demand projected to 2-3x by 2030 — driving grid equipment orders",
        "Transformer lead times 80-120 weeks; transmission queue 2+ years in most markets",
        "EU, US, and China all have multi-year grid investment programs",
        "Policy support: IRA, EU Grid Action Plan, China State Grid Five-Year Plan",
    ]
    add_content_slide(prs, "Electricity Continuity Is Becoming Strategic Infrastructure", industry_items)

    # -------- Slide 5: Why Now / Why History Is Misleading --------
    why_now_items = [
        "⚠️ 2Y backtest can be unfavorable because broad energy-transition winners often rerated together",
        "This is NOT a historically validated spread — it is a forward-looking variant view on regime shift",
        "Past regime (2021-2024): Broad clean-tech beta lifted inverter and equipment winners together",
        "New regime (2025-2030): Grid infrastructure durability beats high-growth clean-tech valuation risk",
        "Stock-pool compliance: Dongfang is the pool anchor; Sungrow is the non-pool same-sector short",
        "We are betting on the BREAKDOWN of historical correlation, not its continuation",
    ]
    add_content_slide(prs, "Why Now: The Regime Shift", why_now_items)

    # -------- Slide 6: Long Case - Dongfang --------
    long_items = [
        "Integrated power equipment and grid-integration leader with clean-energy equipment exposure",
        "Strong domestic grid orders from State Grid Corp + overseas expansion (Middle East, SE Asia)",
        "Revenue growth trajectory supported by multi-year order backlog",
        "Gross margins stable/improving vs peers; R&D investment in digital grid",
        "Variant gap: Market prices it as cyclical China equipment — we price it as global grid-resilience compounder",
    ]
    add_content_slide(prs, SLIDE_TITLES["long_case"], long_items)

    if include_per_slide_evidence and evidence_pack and "slide_5_long_case" in evidence_pack:
        add_quote_slide(
            prs,
            "Long Case Evidence",
            evidence_pack["slide_5_long_case"]["quotes"],
        )

    # -------- Slide 7: Short Case - Sungrow --------
    short_items = [
        "High-expectation inverter and storage leader facing demand normalization",
        "2025 strong results (Revenue RMB 89.2bn, +14.6%; Net Profit RMB 13.5bn, +22.0%)",
        "Q1 2026 shows clear slowdown: Revenue -18.3% YoY, Net Profit -40.1% YoY",
        "Premium multiple vulnerable to inverter/storage price compression and margin pressure",
        "Short thesis is valuation de-rating as growth expectations normalize, not business failure",
    ]
    add_content_slide(prs, SLIDE_TITLES["short_case"], short_items)

    if include_per_slide_evidence and evidence_pack and "slide_7_short_case" in evidence_pack:
        add_quote_slide(
            prs,
            "Short Case Evidence",
            evidence_pack["slide_7_short_case"]["quotes"],
        )

    # -------- Slide 8: Backtest - Short-Leg Correlation --------
    backtest_path = charts_dir / "oil_sungrow_correlation.png"
    if backtest_path.exists():
        add_chart_slide(
            prs,
            f"Historical Divergence: Market vs {LONG_LEG.name} vs {SHORT_LEG.name} (2Y)",
            backtest_path,
            "Historical result is a risk check, not proof. Forward case depends on grid infrastructure durability outperforming high-growth clean-tech valuation risk.",
        )

    scorecard_path = Path("data/processed/valuation/predictive_scorecard.csv")
    if scorecard_path.exists():
        try:
            scorecard_df = pd.read_csv(scorecard_path)
            if not scorecard_df.empty:
                add_table_slide(
                    prs,
                    "Empirical Setup: Forward Predictive Scorecard",
                    scorecard_df[["pillar", "evidence", "signal"]],
                    "Separates adverse historical backtest from current empirical setup and forward catalysts.",
                )
        except Exception:
            pass

    # -------- Slide 9: Comparison Matrix (Chart) --------
    matrix_path = charts_dir / "long_short_matrix.png"
    if matrix_path.exists():
        add_chart_slide(
            prs,
            SLIDE_TITLES["comparison"],
            matrix_path,
            f"Long {LONG_LEG.name} captures grid capex tailwind. Short {SHORT_LEG.name} carries high-expectation inverter and storage de-rating risk.",
        )

    # -------- Slide 10: AI Signal Tracker --------
    if signal_tracker_df is not None and not signal_tracker_df.empty:
        caption = ""
        if narrative_shift:
            caption = (
                f"Thesis support score: {narrative_shift['thesis_support_score']} | "
                f"{narrative_shift['interpretation']} | "
                f"Grid signal share: {narrative_shift['grid_signal_share']*100:.0f}%"
            )
        add_table_slide(
            prs,
            "Thematic Evidence Tracker: Text-Based Signal Compilation",
            signal_tracker_df,
            caption,
        )

    # Detailed text evidence and signal-frequency charts are kept in the
    # generated evidence pack, not the <=20-page submission deck.

    # -------- Slide 11: Peer Comparables --------
    if peer_comps_df is not None and not peer_comps_df.empty:
        add_table_slide(
            prs,
            "Valuation: Peer Comparables",
            peer_comps_df,
            "Comps compare clean grid exposure against broader power-equipment and industrial equipment exposure",
        )

    # -------- Slide 12: DCF Cross-Check --------
    if dcf_df is not None and not dcf_df.empty:
        add_table_slide(
            prs,
            "Valuation: DCF Cross-Check",
            dcf_df,
            "Simple normalized FCF DCF used as a consistency check against the scenario framework",
        )

    # -------- Slide 13: Long Scenarios --------
    if long_scenarios_df is not None and not long_scenarios_df.empty:
        add_table_slide(
            prs,
            "Long Scenarios: Dongfang Electric",
            long_scenarios_df,
            "Probability-weighted upside driven by overseas revenue mix and grid re-rating",
        )

    # -------- Slide 14: Short Scenarios --------
    if short_scenarios_df is not None and not short_scenarios_df.empty:
        add_table_slide(
            prs,
            f"Short Scenarios: {SHORT_LEG.name}",
            short_scenarios_df,
            "Probability-weighted downside driven by thin margins and multiple compression",
        )

    # Sensitivity analysis remains available in outputs/charts and the report;
    # omit it from the main deck to stay within the 20-page limit.

    # -------- Slide 15: Catalysts --------
    catalyst_items = [
        "LONG CATALYSTS (Dongfang):",
        "  • Quarterly overseas order wins (Middle East, SE Asia)",
        "  • China State Grid capex announcements",
        "  • Gross margin expansion on overseas mix shift",
        f"SHORT CATALYSTS ({SHORT_LEG.name}):",
        "  • Margin recovery falls short of expectations",
        "  • Inverter/storage pricing pressure persists",
        "  • Investors de-rate premium clean-tech expectations vs backlog-backed grid exposure",
    ]
    add_content_slide(prs, "Catalysts Over Next 6-12 Months", catalyst_items)

    # -------- Slide 16: Risks --------
    risk_items = [
        f"Broad power-equipment rally may lift {SHORT_LEG.name} temporarily → position sizing limits risk",
        "Dongfang China beta/multiple compression → offset by 140B RMB backlog visibility",
        "Grid capex delays → multi-year policy demand reduces single-year risk",
        "FX risk on RMB → monitored; sensitivity table in appendix",
        f"SHORT SQUEEZE / BORROW COST: {SHORT_LEG.name} borrow cost assumed ~2.5-7%; position sizing limits risk",
        f"Thesis kill switch: {SHORT_LEG.name} margin expansion + Dongfang backlog decline → exit pair",
    ]
    add_content_slide(prs, "Risks and Mitigants", risk_items)

    # -------- Slide 17: Trader Execution Framework --------
    trader_path = Path("data/processed/valuation/trader_analysis.csv")
    trader_row = {}
    if trader_path.exists():
        try:
            trader_df = pd.read_csv(trader_path)
            if not trader_df.empty:
                trader_row = trader_df.iloc[0].to_dict()
        except Exception:
            trader_row = {}
    pair_vol = trader_row.get("volatilities_pair_vol", 83.2)
    rec_notional = trader_row.get("position_sizing_recommended_notional_mm", 1.2)
    long_notional = trader_row.get("position_sizing_long_dongfang_notional_mm", 0.6)
    short_notional = trader_row.get("position_sizing_short_sungrow_notional_mm", 0.6)
    trader_items = [
        "POSITION SIZING (Risk-Based):",
        "  • Portfolio: $100M example | Max risk per trade: 2% ($2M)",
        f"  • Pair volatility: {pair_vol:.1f}% annual | Position size: {rec_notional:.1f}% of portfolio (${rec_notional:.1f}M notional)",
        f"  • Allocation: ${long_notional:.1f}M long Dongfang / ${short_notional:.1f}M short {SHORT_LEG.name} (dollar-neutral)",
        "",
        "CARRY COST (6-month hold):",
        f"  • {SHORT_LEG.name} borrow cost: 2.5-7% annually",
        f"  • {SHORT_LEG.name} dividend/carry checked before execution",
        "  • Dongfang dividend (long receives): 1.5% → ~$9K income",
        "  • Net carry: ~$18K-36K (1.5-3% of expected spread return)",
        "",
        "LIQUIDITY & EXECUTION:",
        "  • Dongfang (HK): ~$147M daily volume | Execute in <1 day | HKEX access",
        f"  • {SHORT_LEG.name}: liquidity checked via current trader_analysis.csv",
        "  • Dongfang trades on HKEX (1072.HK) - no Stock Connect needed",
        "",
        "TECHNICAL TIMING:",
        f"  • {SHORT_LEG.name}: monitor for short entry confirmation",
        "  • Dongfang: Near highs but policy tailwinds support entry",
    ]
    add_content_slide(prs, "Execution: Position Sizing, Carry & Liquidity", trader_items)

    # -------- Slide 18: Text Analysis Module: Honest Assessment --------
    doc_count = 47
    paragraph_count = len(classified_df) if classified_df is not None else 426
    doc_index_path = Path("data/processed/document_index.csv")
    if doc_index_path.exists():
        try:
            doc_count = len(pd.read_csv(doc_index_path))
        except Exception:
            pass
    if consolidate_ai_slides:
        # Single combined slide - reframed honestly
        ai_combined = [
            "WHAT THIS MODULE DOES:",
            f"  • Compiles text evidence from {doc_count} documents / {paragraph_count} paragraphs into structured categories",
            "  • Surfaces specific quotes and management language for manual review",
            "  • Feeds a predictive scorecard with fundamentals, valuation, technicals, and catalysts",
            "CLAIM BOUNDARY:",
            "  • Predictive input, not autonomous forecast — text signals do not trade themselves",
            "  • Corpus is useful for thesis discovery, but still too small for standalone statistical alpha claims",
            "  • Classification relies on keyword patterns, not deep semantic understanding",
            "  • Source bias: RSS feeds overweight news vs. operational 8-K filings",
            "  • Every output requires human verification before investment decisions",
        ]
        add_content_slide(prs, "Text Analysis: Scope and Limitations", ai_combined)
    else:
        # Two separate slides - reframed
        add_content_slide(prs, "Text Analysis: What It Provides", [
            f"Compiles text evidence from {doc_count} documents / {paragraph_count} paragraphs into structured categories",
            "Surfaces specific quotes and management language for manual review",
            "Feeds a predictive scorecard with fundamentals, valuation, technicals, and catalysts",
            "Full source traceability — every quote links to original document",
        ])
        add_content_slide(prs, "Text Analysis: Honest Limitations", [
            "Predictive input, not autonomous forecast — text signals do not trade themselves",
            f"{paragraph_count} paragraphs = limited sample for statistical return forecasting",
            "Classification relies on keyword patterns, not deep semantic analysis",
            "Source bias: RSS feeds overweight news vs. operational filings",
            "Every output requires human verification before investment use",
        ])

    # Duplicate money/backtest slides are omitted from the main deck; their
    # content is already covered in the executive summary, scorecard, and risk slides.

    # -------- Slide 19: Recommendation --------
    rec_subtitle = (
        f"Long {LONG_LEG.name} / Short {SHORT_LEG.name}"
        if not valuation_summary else
        f"Long Dongfang ({valuation_summary['long_expected_return']:+.0f}%) / "
        f"Short {SHORT_LEG.name} ({valuation_summary['short_expected_return']:+.0f}%) = "
        f"Pair spread {valuation_summary['pair_spread_return']:+.0f}%"
    )
    add_title_slide(prs, "Recommendation", rec_subtitle)

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"[SAVED] PPTX: {output_path}")
    return output_path


def main():
    """Build the full deck from processed outputs."""
    from src.config import (
        CLASSIFIED_PARAGRAPHS_PATH,
        CHARTS_DIR,
        OUTPUTS_DIR,
        PROCESSED_DIR,
    )
    from src.analysis import plan_format_signal_tracker, narrative_shift_analysis

    # Load classifications
    classified_df = None
    if CLASSIFIED_PARAGRAPHS_PATH.exists():
        try:
            classified_df = pd.read_csv(CLASSIFIED_PARAGRAPHS_PATH)
            print(f"Loaded {len(classified_df)} classified paragraphs")
        except pd.errors.EmptyDataError:
            print("Classifications file empty")

    # Build signal tracker
    signal_tracker_df = None
    narrative = None
    if classified_df is not None and not classified_df.empty:
        signal_tracker_df = plan_format_signal_tracker(classified_df)
        narrative = narrative_shift_analysis(classified_df)

    # Load evidence pack
    evidence_pack = None
    evidence_json = OUTPUTS_DIR / "tables" / "evidence_pack.json"
    if evidence_json.exists():
        with open(evidence_json) as f:
            evidence_pack = json.load(f)

    # Load valuation outputs
    val_dir = PROCESSED_DIR / "valuation"
    peer_comps_df = None
    dcf_df = None
    long_scenarios_df = None
    short_scenarios_df = None
    valuation_summary = None

    if (val_dir / "peer_comps.csv").exists():
        peer_comps_df = pd.read_csv(val_dir / "peer_comps.csv")
    if (val_dir / "dcf_cross_check.csv").exists():
        dcf_df = pd.read_csv(val_dir / "dcf_cross_check.csv")
    if (val_dir / "long_scenarios.csv").exists():
        long_scenarios_df = pd.read_csv(val_dir / "long_scenarios.csv")
    if (val_dir / "short_scenarios.csv").exists():
        short_scenarios_df = pd.read_csv(val_dir / "short_scenarios.csv")
    if (val_dir / "pair_trade_summary.csv").exists():
        pair = pd.read_csv(val_dir / "pair_trade_summary.csv").iloc[0].to_dict()
        valuation_summary = {
            "long_expected_return": pair.get("long_expected_return_pct", 0),
            "short_expected_return": pair.get("short_expected_move_pct", 0),
            "pair_spread_return": pair.get("pair_spread_return_pct", 0),
        }

    # Build deck
    deck_dir = Path(__file__).parent.parent / "deck"
    deck_dir.mkdir(parents=True, exist_ok=True)
    output_path = deck_dir / "UBS_Pitch_Deck_AUTO.pptx"

    build_deck(
        classified_df=classified_df,
        signal_tracker_df=signal_tracker_df,
        evidence_pack=evidence_pack,
        valuation_summary=valuation_summary,
        peer_comps_df=peer_comps_df,
        dcf_df=dcf_df,
        long_scenarios_df=long_scenarios_df,
        short_scenarios_df=short_scenarios_df,
        charts_dir=CHARTS_DIR,
        output_path=output_path,
        narrative_shift=narrative,
    )

    print("\n" + "=" * 60)
    print("DECK GENERATED")
    print("=" * 60)
    print(f"Open: {output_path}")


if __name__ == "__main__":
    main()
